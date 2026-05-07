from __future__ import annotations

from datetime import datetime
from pathlib import Path
from typing import Any
import csv
import hashlib
import json
import logging
import math
import re
import shutil
import httpx

from llama_index.core import Document
from llama_index.core.embeddings import BaseEmbedding
from llama_index.core.node_parser import MarkdownNodeParser, SentenceSplitter, TokenTextSplitter
from llama_index.core.schema import TextNode
from sqlmodel import or_, select

from app.core.config import settings
from app.core.database import async_session_factory
from app.models import KnowledgeBase, Model, Provider
from app.integrations.llm.client import llm_embed


logger = logging.getLogger(__name__)

_TOKEN_RE = re.compile(r"\w+", re.UNICODE)
_HTML_RE = re.compile(r"<[^>]+>")
_PLACEHOLDER_RE = re.compile(r"\{\{\s*([\w.]+)\s*\}\}")
_STORED_UPLOAD_PREFIX_RE = re.compile(r"^[0-9a-fA-F-]{36}_(.+)$")
_EMBED_DIM = 256


def _safe_json_loads(value: str | None, fallback: Any):
    if not value:
        return fallback
    try:
        return json.loads(value)
    except Exception:
        return fallback


def _kb_storage_dir(kb_id: str) -> Path:
    path = settings.data_dir / "vector_store" / kb_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def clear_kb_storage(kb_id: str) -> None:
    storage_dir = settings.data_dir / "vector_store" / kb_id
    if storage_dir.exists():
        shutil.rmtree(storage_dir, ignore_errors=True)


def _tokenize(text: str) -> list[str]:
    words = _TOKEN_RE.findall(text.lower())
    chars = [char for char in text if not char.isspace()]
    bigrams = ["".join(chars[index:index + 2]) for index in range(len(chars) - 1)]
    return words + bigrams


def _hashed_embedding(text: str, dims: int = _EMBED_DIM) -> list[float]:
    vector = [0.0] * dims
    for token in _tokenize(text):
        digest = hashlib.sha256(token.encode("utf-8")).digest()
        bucket = int.from_bytes(digest[:4], "big") % dims
        sign = 1.0 if digest[4] % 2 == 0 else -1.0
        vector[bucket] += sign
    norm = math.sqrt(sum(value * value for value in vector)) or 1.0
    return [round(value / norm, 6) for value in vector]


def _keyword_overlap_score(query: str, text: str) -> float:
    query_tokens = set(_tokenize(query))
    if not query_tokens:
        return 0.0
    text_tokens = set(_tokenize(text))
    return len(query_tokens & text_tokens) / len(query_tokens)


def _tokenize_words_only(text: str) -> list[str]:
    normalized = text.lower().replace("_", " ")
    return re.findall(r"[0-9a-z\u4e00-\u9fff]+", normalized)


def _lexical_rerank_score(query: str, text: str) -> float:
    normalized_query = " ".join(query.lower().split())
    normalized_text = " ".join(text.lower().split())
    exact_phrase_score = 1.0 if normalized_query and normalized_query in normalized_text else 0.0

    query_terms = set(_tokenize_words_only(query))
    if not query_terms:
        return exact_phrase_score

    text_terms = set(_tokenize_words_only(text))
    term_overlap = len(query_terms & text_terms) / len(query_terms)
    if exact_phrase_score > 0.0:
        return 1.0
    return round(term_overlap * 0.4, 4)


def _embedding_url(api_base: str | None) -> str:
    base = (api_base or "https://api.openai.com/v1").rstrip("/")
    if base.endswith("/v1"):
        return f"{base}/embeddings"
    return f"{base}/v1/embeddings"


class ProviderEmbedding(BaseEmbedding):
    model_name: str
    api_key: str | None = None
    api_base: str | None = None

    @classmethod
    def class_name(cls) -> str:
        return "ProviderEmbedding"

    def _embed_via_provider(self, text: str) -> list[float]:
        if self.model_name == "local/hash-embedding":
            return _hashed_embedding(text)
        try:
            payload = {
                "model": self.model_name,
                "input": [text],
            }
            headers = {"Content-Type": "application/json"}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            with httpx.Client(timeout=30) as client:
                response = client.post(
                    _embedding_url(self.api_base),
                    headers=headers,
                    json=payload,
                )
                response.raise_for_status()
                data = response.json()
            rows = data.get("data") or []
            if rows and isinstance(rows[0].get("embedding"), list):
                return [float(value) for value in rows[0]["embedding"]]
        except Exception as exc:
            logger.warning("Embedding fallback for %s: %s", self.model_name, exc)
        return _hashed_embedding(text)

    def _get_query_embedding(self, query: str) -> list[float]:
        return self._embed_via_provider(query)

    async def _aget_query_embedding(self, query: str) -> list[float]:
        if self.model_name == "local/hash-embedding":
            return _hashed_embedding(query)
        try:
            return await llm_embed(
                model=self.model_name,
                text=query,
                api_key=self.api_key,
                api_base=self.api_base,
            )
        except Exception as exc:
            logger.warning("Async embedding fallback for %s: %s", self.model_name, exc)
            return _hashed_embedding(query)

    def _get_text_embedding(self, text: str) -> list[float]:
        return self._embed_via_provider(text)


async def resolve_kb_runtime(kb_id: str) -> dict[str, Any]:
    async with async_session_factory() as session:
        kb = await session.get(KnowledgeBase, kb_id)
        if not kb:
            raise ValueError(f"Knowledge base {kb_id} not found")

        model = None
        provider = None
        if kb.embedding_model_id:
            result = await session.exec(
                select(Model).where(
                    or_(Model.id == kb.embedding_model_id, Model.model_id == kb.embedding_model_id)
                )
            )
            model = result.first()
        if model:
            provider = await session.get(Provider, model.provider_id)
        elif kb.embedding_provider_type:
            provider_result = await session.exec(
                select(Provider).where(
                    Provider.provider_type == kb.embedding_provider_type,
                    Provider.enabled == True,
                )
            )
            provider = provider_result.first()

        model_name = model.model_id if model else (kb.embedding_model_id or "local/hash-embedding")
        return {
            "kb": kb,
            "provider": provider,
            "model": model,
            "embed_model": ProviderEmbedding(
                model_name=model_name,
                api_key=provider.api_key_encrypted if provider else None,
                api_base=provider.base_url if provider else None,
            ),
            "persist_dir": _kb_storage_dir(kb_id),
            "parser_config": _safe_json_loads(kb.parser_config_json, {}),
            "retrieval_config": _safe_json_loads(kb.retrieval_config_json, {}),
            "metadata_template": _safe_json_loads(kb.metadata_template_json, {}),
        }


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _read_pdf(path: Path) -> str:
    import pypdf

    reader = pypdf.PdfReader(str(path))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _read_docx(path: Path) -> str:
    from docx import Document as DocxDocument

    document = DocxDocument(str(path))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for sheet in workbook.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None and str(value).strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"# Slide {slide_index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _read_csv(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as file:
        reader = csv.reader(file)
        rows = [" | ".join(row) for row in reader]
    return "\n".join(rows)


def _read_json(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as file:
        payload = json.load(file)
    return json.dumps(payload, ensure_ascii=False, indent=2)


def _read_html(path: Path) -> str:
    raw = _read_text_file(path)
    return re.sub(r"\s+", " ", _HTML_RE.sub(" ", raw)).strip()


def _read_with_markitdown(file_path: str, parser_config: dict[str, Any] | None = None) -> dict[str, Any]:
    parser_config = parser_config or {}
    markitdown_config = parser_config.get("markitdown") or {}
    fallback_to_native = markitdown_config.get("fallback_to_native", True)
    prefer_markdown_chunking = markitdown_config.get("prefer_markdown_chunking", True)
    enable_plugins = markitdown_config.get("use_plugins", False)
    try:
        from markitdown import MarkItDown

        converter = MarkItDown(enable_plugins=enable_plugins)
        result = converter.convert(file_path)
        text = getattr(result, "text_content", None) or ""
        if not text.strip():
            raise ValueError("MarkItDown returned empty content")
        return {
            "text": text,
            "preprocessor": "markitdown",
            "chunk_strategy_override": "markdown" if prefer_markdown_chunking else None,
        }
    except Exception as exc:
        if fallback_to_native:
            logger.warning("MarkItDown preprocessing failed for %s: %s", file_path, exc)
            return {
                "text": read_source_text(file_path),
                "preprocessor": "native-fallback",
                "chunk_strategy_override": None,
            }
        raise


def read_source_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _read_pdf(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".xlsx":
        return _read_xlsx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix == ".csv":
        return _read_csv(path)
    if suffix == ".tsv":
        return _read_text_file(path).replace("\t", " | ")
    if suffix == ".json":
        return _read_json(path)
    if suffix in {".html", ".htm"}:
        return _read_html(path)
    return _read_text_file(path)


def read_source_content(file_path: str, parser_config: dict[str, Any] | None = None) -> dict[str, Any]:
    parser_config = parser_config or {}
    preprocessor = str(parser_config.get("preprocessor") or "native").lower()
    if preprocessor == "markitdown":
        return _read_with_markitdown(file_path, parser_config)
    return {
        "text": read_source_text(file_path),
        "preprocessor": "native",
        "chunk_strategy_override": None,
    }


def _resolve_placeholder(values: dict[str, Any], key: str):
    current: Any = values
    for part in key.split("."):
        if isinstance(current, dict) and part in current:
            current = current[part]
        else:
            return None
    return current


def _render_template_value(value: Any, variables: dict[str, Any]) -> Any:
    if isinstance(value, dict):
        return {k: _render_template_value(v, variables) for k, v in value.items()}
    if isinstance(value, list):
        return [_render_template_value(v, variables) for v in value]
    if not isinstance(value, str):
        return value

    matches = list(_PLACEHOLDER_RE.finditer(value))
    if not matches:
        return value
    if len(matches) == 1 and matches[0].span() == (0, len(value)):
        resolved = _resolve_placeholder(variables, matches[0].group(1))
        return resolved if resolved is not None else value

    def _replace(match: re.Match[str]) -> str:
        resolved = _resolve_placeholder(variables, match.group(1))
        if resolved is None:
            return match.group(0)
        if isinstance(resolved, (dict, list)):
            return json.dumps(resolved, ensure_ascii=False)
        return str(resolved)

    return _PLACEHOLDER_RE.sub(_replace, value)


def build_metadata(
    kb_id: str,
    file_path: str,
    metadata_mode: str,
    metadata_template: dict[str, Any] | None,
    source_filename: str | None = None,
) -> dict[str, Any]:
    path = Path(file_path)
    resolved_filename = source_filename or path.name
    variables = {
        "kb_id": kb_id,
        "filename": resolved_filename,
        "file_path": str(path),
        "extension": path.suffix.lower().lstrip("."),
        "timestamp": datetime.utcnow().isoformat(),
    }
    metadata = {
        "kb_id": kb_id,
        "filename": resolved_filename,
        "source": str(path),
        "extension": path.suffix.lower().lstrip("."),
        "imported_at": variables["timestamp"],
    }
    if metadata_mode == "disabled":
        return {"kb_id": kb_id, "source": str(path), "filename": resolved_filename}
    if metadata_mode == "custom" and metadata_template:
        metadata.update(_render_template_value(metadata_template, variables))
    return metadata


def _paragraph_nodes(text: str, metadata: dict[str, Any], chunk_size: int, chunk_overlap: int) -> list[TextNode]:
    splitter = SentenceSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    nodes: list[TextNode] = []
    chunk_index = 0
    for paragraph in re.split(r"\n\s*\n", text):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        if len(paragraph) <= chunk_size:
            nodes.append(TextNode(text=paragraph, metadata={**metadata, "chunk_index": chunk_index}))
            chunk_index += 1
            continue
        paragraph_doc = Document(text=paragraph, metadata=metadata)
        for node in splitter.get_nodes_from_documents([paragraph_doc]):
            node.metadata = {**node.metadata, "chunk_index": chunk_index}
            nodes.append(node)
            chunk_index += 1
    return nodes


def build_nodes(
    text: str,
    metadata: dict[str, Any],
    chunk_strategy: str,
    chunk_size: int,
    chunk_overlap: int,
    parser_config: dict[str, Any] | None = None,
) -> list[TextNode]:
    parser_config = parser_config or {}
    resolved_chunk_size = int(parser_config.get("chunk_size") or chunk_size)
    resolved_chunk_overlap = int(parser_config.get("chunk_overlap") or chunk_overlap)
    document = Document(text=text, metadata=metadata)

    if chunk_strategy == "markdown":
        parser = MarkdownNodeParser.from_defaults(
            header_path_separator=parser_config.get("header_path_separator", "/")
        )
        nodes = parser.get_nodes_from_documents([document])
    elif chunk_strategy == "token":
        parser = TokenTextSplitter(
            chunk_size=resolved_chunk_size,
            chunk_overlap=resolved_chunk_overlap,
        )
        nodes = parser.get_nodes_from_documents([document])
    elif chunk_strategy == "paragraph":
        nodes = _paragraph_nodes(text, metadata, resolved_chunk_size, resolved_chunk_overlap)
    else:
        parser = SentenceSplitter(
            chunk_size=resolved_chunk_size,
            chunk_overlap=resolved_chunk_overlap,
        )
        nodes = parser.get_nodes_from_documents([document])

    normalized_nodes: list[TextNode] = []
    for index, node in enumerate(nodes):
        node.metadata = {**metadata, **node.metadata, "chunk_index": index}
        normalized_nodes.append(node)
    return normalized_nodes


def _get_lance_dir(kb_id: str) -> Path:
    path = settings.data_dir / "vector_store" / kb_id / "lance"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _escape_sql(value: str) -> str:
    return value.replace("'", "''")


def _normalize_metadata_filter_value(value: Any) -> Any:
    if isinstance(value, str):
        return value.strip().casefold()
    return value


def _candidate_metadata_filter_values(key: str, value: Any) -> set[Any]:
    normalized_value = _normalize_metadata_filter_value(value)
    candidates = {normalized_value}
    if key == "filename" and isinstance(value, str):
        match = _STORED_UPLOAD_PREFIX_RE.match(value.strip())
        if match:
            candidates.add(match.group(1).casefold())
    return candidates


def get_document_chunk_count(kb_id: str, source_path: str) -> int:
    try:
        import lancedb
        db = lancedb.connect(str(_get_lance_dir(kb_id)))
        if "chunks" not in db.table_names():
            return 0
        table = db.open_table("chunks")
        return table.count_rows(filter=f"source = '{_escape_sql(source_path)}'")
    except Exception:
        return 0


def delete_document_nodes(kb_id: str, source_path: str) -> int:
    """Remove all vector nodes for a specific document without rebuilding the whole index."""
    try:
        import lancedb
        db = lancedb.connect(str(_get_lance_dir(kb_id)))
        if "chunks" not in db.table_names():
            return 0
        table = db.open_table("chunks")
        count_before = table.count_rows()
        table.delete(f"source = '{_escape_sql(source_path)}'")
        count_after = table.count_rows()
        return count_before - count_after
    except Exception as exc:
        logger.warning("Failed to delete nodes for kb %s: %s", kb_id, exc)
        return 0


def add_nodes_to_vector_store(kb_id: str, nodes: list[TextNode], embed_model: BaseEmbedding) -> int:
    import lancedb
    records = []
    for node in nodes:
        text = node.get_content()
        embedding = embed_model.get_text_embedding(text)
        metadata = node.metadata or {}
        records.append({
            "node_id": node.node_id,
            "text": text,
            "vector": list(embedding),
            "source": metadata.get("source", ""),
            "filename": metadata.get("filename", ""),
            "metadata_json": json.dumps(metadata, ensure_ascii=False),
        })
    if not records:
        return 0
    db = lancedb.connect(str(_get_lance_dir(kb_id)))
    if "chunks" in db.table_names():
        table = db.open_table("chunks")
        table.add(records)
    else:
        db.create_table("chunks", data=records)
    return len(records)


async def query_vector_store(
    kb_id: str,
    query: str,
    top_k: int,
    embed_model: BaseEmbedding,
) -> list[dict[str, Any]]:
    import lancedb
    db = lancedb.connect(str(_get_lance_dir(kb_id)))
    if "chunks" not in db.table_names():
        return []
    table = db.open_table("chunks")
    query_embedding = await embed_model.aget_query_embedding(query)
    try:
        rows = (
            table.search(list(query_embedding), vector_column_name="vector")
            .metric("cosine")
            .limit(top_k)
            .to_list()
        )
    except Exception as exc:
        logger.warning("Vector search failed for kb %s: %s", kb_id, exc)
        return []
    results: list[dict[str, Any]] = []
    for row in rows:
        metadata = json.loads(row.get("metadata_json") or "{}")
        distance = float(row.get("_distance", 1.0))
        score = round(1.0 - distance, 4)
        results.append({
            "text": row.get("text", ""),
            "score": score,
            "source": row.get("source", "unknown"),
            "metadata": metadata,
            "node_id": row.get("node_id", ""),
        })
    return results


def apply_metadata_filters(results: list[dict[str, Any]], metadata_filters: dict[str, Any] | None) -> list[dict[str, Any]]:
    if not metadata_filters:
        return results
    filtered = []
    for item in results:
        metadata = item.get("metadata") or {}
        if all(
            _normalize_metadata_filter_value(value) in _candidate_metadata_filter_values(key, metadata.get(key))
            for key, value in metadata_filters.items()
        ):
            filtered.append(item)
    return filtered


def apply_score_threshold(results: list[dict[str, Any]], min_score: float | None) -> list[dict[str, Any]]:
    if min_score is None:
        return results
    return [
        item for item in results if float(item.get("score") or 0.0) >= min_score
    ]


def hybrid_rerank(results: list[dict[str, Any]], query: str) -> list[dict[str, Any]]:
    reranked = []
    for item in results:
        lexical_score = _lexical_rerank_score(query, item.get("text", ""))
        base_score = float(item.get("score") or 0.0)
        final_score = round(base_score * 0.65 + lexical_score * 0.35, 4)
        reranked.append({**item, "score": final_score, "lexical_score": round(lexical_score, 4)})
    reranked.sort(key=lambda item: item.get("score", 0.0), reverse=True)
    return reranked