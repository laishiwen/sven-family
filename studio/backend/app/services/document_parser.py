"""Document parser: PDF, Office (docx, xlsx, pptx), images (OCR), plain text."""

from __future__ import annotations

import base64
import io
import logging
import os
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


async def parse_document(file_path: str | Path, mime_type: str | None = None) -> dict[str, Any]:
    """Parse a document file and return structured content."""
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {file_path}")

    ext = path.suffix.lower()
    if not mime_type:
        mime_type = _guess_mime(ext)

    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            content = await _parse_pdf(path)
        elif ext in (".docx", ".doc") or mime_type and "word" in mime_type:
            content = await _parse_docx(path)
        elif ext in (".xlsx", ".xls") or mime_type and "excel" in mime_type:
            content = await _parse_xlsx(path)
        elif ext in (".pptx", ".ppt") or mime_type and "presentation" in mime_type:
            content = await _parse_pptx(path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            content = await _parse_image(path)
        elif ext in (".txt", ".md", ".json", ".csv", ".xml", ".yaml", ".yml", ".py", ".js", ".ts", ".html", ".css"):
            content = path.read_text(encoding="utf-8", errors="replace")
        else:
            content = path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logger.warning("Failed to parse %s with %s: %s, falling back to raw text", file_path, mime_type, e)
        try:
            content = path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            content = f"[Unable to parse document: {e}]"

    return {
        "filename": path.name,
        "mime_type": mime_type,
        "content": content,
        "size_bytes": path.stat().st_size,
    }


def _guess_mime(ext: str) -> str:
    mapping = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc": "application/msword",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".json": "application/json",
        ".csv": "text/csv",
    }
    return mapping.get(ext, "application/octet-stream")


async def _parse_pdf(path: Path) -> str:
    """Extract text from PDF."""
    try:
        import pymupdf
        doc = pymupdf.open(str(path))
        pages = []
        for page in doc:
            pages.append(page.get_text("text"))
        doc.close()
        return "\n\n".join(pages)
    except ImportError:
        pass

    try:
        from pdfplumber import open as pdf_open
        with pdf_open(str(path)) as pdf:
            pages = [page.extract_text() or "" for page in pdf.pages]
        return "\n\n".join(pages)
    except ImportError:
        pass

    raise ImportError("Install pymupdf or pdfplumber for PDF parsing: pip install pymupdf")


async def _parse_docx(path: Path) -> str:
    """Extract text from DOCX."""
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except ImportError:
        raise ImportError("Install python-docx for DOCX parsing: pip install python-docx")


async def _parse_xlsx(path: Path) -> str:
    """Extract text from XLSX."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append("\t".join(str(c) if c is not None else "" for c in row))
            parts.append("\n".join(rows[:1000]))  # Cap at 1k rows
        wb.close()
        return "\n\n".join(parts)
    except ImportError:
        raise ImportError("Install openpyxl for XLSX parsing: pip install openpyxl")


async def _parse_pptx(path: Path) -> str:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            if texts:
                slides.append(f"## Slide {i}\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise ImportError("Install python-pptx for PPTX parsing: pip install python-pptx")


async def _parse_image(path: Path) -> str:
    """Extract text from image using OCR."""
    # Try Tesseract OCR first
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        if text.strip():
            return f"[OCR Result]\n{text}"
    except ImportError:
        pass

    # Fallback: describe the image
    return f"[Image: {path.name}] (OCR not available — install pytesseract and tesseract-ocr)"


async def parse_bytes(data: bytes, filename: str, mime_type: str | None = None) -> dict[str, Any]:
    """Parse document from bytes."""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as f:
        f.write(data)
        tmp_path = f.name
    try:
        return await parse_document(tmp_path, mime_type)
    finally:
        os.unlink(tmp_path)
